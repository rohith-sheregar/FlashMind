"""
Robust text extraction, cleaning and chunking utilities for academic documents.

Requires:
# pip install PyMuPDF python-docx python-pptx chardet

This module provides a clean API to extract text from PDF/DOCX/PPTX/TXT,
clean the text for prompt-based workflows, detect headings, and produce
meaningful chunks with optional overlap.
"""

from pathlib import Path
import re
import logging
from typing import List, Dict, Optional

logger = logging.getLogger(__name__)

_HEADER_PATTERNS = [
    r"^\s*page\s+\d+\b",
    r"^lecture notes\b",
    r"^lecture\b",
    r"^slide\b",
    r"^session\b",
    r"^table of contents\b",
    r"\bauthor\b",
    r"\bisbn\b",
    r"\bdoi\b",
    r"©",
    r"all rights reserved",
    r"^[A-Z]{2,}\d{3,}\b",
]


def extract_text_from_file(path: str) -> str:
    """
    Extract and return the full cleaned text from a file.

    Detects file type by suffix and calls the appropriate extractor.

    Args:
        path: path to the input file (.pdf, .docx, .pptx, .txt)

    Returns:
        Cleaned text string.

    Raises:
        ValueError: for unsupported file types or empty extraction.
    """
    p = Path(path)
    suffix = p.suffix.lower()
    if suffix == ".pdf":
        paragraphs = _extract_pdf_paragraphs(path)
        raw = "\n\n".join([p["text"] for p in paragraphs])
    elif suffix in (".docx",):
        raw = _extract_docx_raw(path)
    elif suffix in (".pptx",):
        raw = _extract_pptx_raw(path)
    elif suffix == ".txt":
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            raw = f.read()
    else:
        raise ValueError(f"Unsupported file type: {suffix}")

    cleaned = clean_text(raw)
    if not cleaned.strip():
        raise ValueError("Extraction produced no usable text. File may be image-based (scanned PDF). Consider enabling OCR.")
    return cleaned


def _extract_pdf_paragraphs(path: str) -> List[Dict]:
    """
    Extract text from PDF using PyMuPDF and return paragraph-level list with page numbers.

    Returns a list of dicts: {"text": paragraph_text, "page": page_number}

    Raises ValueError if PDF appears to be image-only (no extractable text).
    """
    try:
        import fitz  # PyMuPDF
    except Exception as e:
        logger.exception("PyMuPDF (fitz) is required for PDF extraction: %s", e)
        raise

    doc = fitz.open(path)
    all_paras: List[Dict] = []
    total_text_found = 0
    for i, page in enumerate(doc, start=1):
        page_text = page.get_text()
        if page_text and page_text.strip():
            total_text_found += 1
        # Split page into lines then paragraphs by blank lines
        lines = [ln.rstrip() for ln in page_text.splitlines()]
        # Remove header/footer style lines per-line
        cleaned_lines = []
        for line in lines:
            s = line.strip()
            if not s:
                cleaned_lines.append("")
                continue
            skip = False
            for pat in _HEADER_PATTERNS:
                if re.search(pat, s, re.IGNORECASE):
                    skip = True
                    break
            if skip:
                continue
            cleaned_lines.append(s)
        # Build paragraphs on page
        para_buf: List[str] = []
        for line in cleaned_lines:
            if not line:
                if para_buf:
                    para = " ".join(para_buf).strip()
                    if para:
                        all_paras.append({"text": para, "page": i})
                    para_buf = []
            else:
                para_buf.append(line)
        if para_buf:
            para = " ".join(para_buf).strip()
            if para:
                all_paras.append({"text": para, "page": i})
    if total_text_found == 0:
        # No extractable text — likely scanned PDF
        raise ValueError("No extractable text found in PDF. File may be image-based (scanned). Enable OCR (e.g., pytesseract) to extract text.")
    return all_paras


def _extract_docx_raw(path: str) -> str:
    """Extract raw paragraph text from DOCX file.

    Returns a plain string joining paragraphs with double-newlines.
    """
    try:
        from docx import Document
    except Exception as e:
        logger.exception("python-docx is required for DOCX extraction: %s", e)
        raise
    doc = Document(path)
    paras = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    return "\n\n".join(paras)


def _extract_pptx_raw(path: str) -> str:
    """Extract raw text from PPTX slides.

    Returns a plain string joining slide text with double-newlines.
    """
    try:
        from pptx import Presentation
    except Exception as e:
        logger.exception("python-pptx is required for PPTX extraction: %s", e)
        raise
    prs = Presentation(path)
    slide_texts: List[str] = []
    for slide in prs.slides:
        bits: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                bits.append(shape.text.strip())
        if bits:
            slide_texts.append(" ".join(bits))
    return "\n\n".join(slide_texts)


# --- Cleaning utilities ---
def clean_text(raw: str) -> str:
    """
    Clean raw extracted text for downstream NLP/model input.

    Applies header/footer removal, hyphenation fixing, line-join heuristics,
    removes numeric-only lines and known artifacts, and normalizes whitespace.
    """
    if not raw:
        return ""
    text = raw
    # Normalize line endings
    text = text.replace('\r\n', '\n').replace('\r', '\n')

    # Fix hyphenation across line breaks: 'exam-\nple' -> 'example'
    text = re.sub(r"(\w)-\n(\w)", r"\1\2", text)

    # Merge lines where a sentence was broken: if line doesn't end with punctuation and next starts with lowercase
    lines = text.split('\n')
    merged_lines: List[str] = []
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        if not line:
            merged_lines.append("")
            i += 1
            continue
        # skip header-like lines
        skip = False
        for pat in _HEADER_PATTERNS:
            if re.search(pat, line, re.IGNORECASE):
                skip = True
                break
        if skip:
            i += 1
            continue

        # Merge with following short/lowercase-starting lines
        j = i + 1
        buf = [line]
        while j < len(lines):
            nxt = lines[j].strip()
            if not nxt:
                break
            # If current buf ends with punctuation, stop merging
            if re.search(r"[\.!?;:]$", buf[-1]):
                break
            # If next line starts lowercase or is short, merge
            if re.match(r"^[a-z]", nxt) or len(nxt) < 60:
                buf.append(nxt)
                j += 1
            else:
                break
        merged = " ".join(buf)
        merged_lines.append(merged)
        i = j

    # Remove lines that are purely numeric or very short page markers
    cleaned = []
    prev = None
    for ln in merged_lines:
        s = ln.strip()
        if not s:
            cleaned.append("")
            prev = s
            continue
        # numeric-only or tokens like '10', 'Page 5'
        if re.fullmatch(r"\d{1,4}", s) or re.match(r"^page\s+\d+\b", s, re.I):
            continue
        # short repeated lines
        if s == prev:
            continue
        # remove ISBN/DOI lines
        if re.search(r"isbn|doi", s, re.I):
            continue
        # remove copyright lines
        if re.search(r"all rights reserved|©", s, re.I):
            continue
        cleaned.append(s)
        prev = s

    # Collapse multiple blank lines to at most one blank line (paragraph separator)
    out = []
    blank = False
    for ln in cleaned:
        if not ln:
            if not blank:
                out.append("")
            blank = True
        else:
            out.append(ln)
            blank = False

    result = "\n\n".join([p for p in "\n\n".join(out).split("\n\n") if p.strip()])
    # Normalize whitespace inside paragraphs
    result = re.sub(r"[ \t]+", " ", result)
    return result.strip()


def _is_heading_line(line: str) -> Optional[str]:
    """Return normalized heading string if line looks like a heading, otherwise None."""
    if not line or len(line) < 3:
        return None
    # Common heading starts
    m = re.match(r"^(Module|Chapter|Section|Unit|Topic)\b[:\s\-]*(.*)$", line, re.I)
    if m:
        left = m.group(1).title()
        rest = m.group(2).strip()
        heading = f"{left} {rest}" if rest else left
        return heading.strip(" -:\t\n")
    # Heuristic: line is mostly Title Case or ALL CAPS and relatively short
    words = line.split()
    if 1 < len(words) <= 8:
        # all caps or title-like
        if line.isupper() or all(w[0].isupper() for w in words if w):
            return line.strip()
    return None


def detect_headings_and_map(text: str) -> List[Dict]:
    """
    Detect headings in the text and return a list mapping paragraph indices to headings.

    Returns a list of dicts: {"para_index": int, "heading": Optional[str]}
    """
    paras = [p.strip() for p in text.split("\n\n") if p.strip()]
    mapping: List[Dict] = []
    current_heading: Optional[str] = None
    for idx, para in enumerate(paras):
        h = _is_heading_line(para.splitlines()[0])
        if h:
            current_heading = h
            mapping.append({"para_index": idx, "heading": current_heading})
        else:
            mapping.append({"para_index": idx, "heading": current_heading})
    return mapping


def extract_text_chunks(path: str,
                        max_chunk_chars: int = 2000,
                        overlap_chars: int = 200,
                        detect_headings: bool = True) -> List[Dict]:
    """
    Extract cleaned text and split into meaningful chunks.

    Each chunk dict contains: chunk_id, text, heading (or None), page_start, page_end
    """
    p = Path(path)
    suffix = p.suffix.lower()
    paragraphs: List[Dict[str, Optional[int]]] = []  # {'text': str, 'page': Optional[int]}

    if suffix == ".pdf":
        paragraphs = _extract_pdf_paragraphs(path)
    elif suffix in (".docx",):
        raw = _extract_docx_raw(path)
        paras = [para.strip() for para in raw.split("\n\n") if para.strip()]
        paragraphs = [{"text": para, "page": None} for para in paras]
    elif suffix in (".pptx",):
        raw = _extract_pptx_raw(path)
        paras = [para.strip() for para in raw.split("\n\n") if para.strip()]
        paragraphs = [{"text": para, "page": None} for para in paras]
    elif suffix == ".txt":
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            raw = f.read()
        paras = [para.strip() for para in raw.split("\n\n") if para.strip()]
        paragraphs = [{"text": para, "page": None} for para in paras]
    else:
        raise ValueError(f"Unsupported file type for chunking: {suffix}")

    # Clean each paragraph individually
    cleaned_paras: List[Dict[str, Optional[int]]] = []
    for item in paragraphs:
        txt = clean_text(item.get("text", ""))
        if txt:
            cleaned_paras.append({"text": txt, "page": item.get("page")})

    if not cleaned_paras:
        raise ValueError("No extractable text to chunk after cleaning.")

    # Merge short paragraphs when appropriate
    merged_paras: List[Dict[str, Optional[int]]] = []
    i = 0
    while i < len(cleaned_paras):
        cur = cleaned_paras[i]
        cur_text = cur["text"]
        # merge with next if short and next starts lowercase
        if len(cur_text) < 150 and i + 1 < len(cleaned_paras):
            nxt = cleaned_paras[i + 1]["text"]
            if re.match(r"^[a-z]", nxt):
                merged_text = cur_text + " " + nxt
                merged_paras.append({"text": merged_text, "page": cur.get("page") or cleaned_paras[i + 1].get("page")})
                i += 2
                continue
        merged_paras.append(cur)
        i += 1

    # Assign headings by scanning merged_paras if requested
    headings_map: List[Optional[str]] = []
    if detect_headings:
        # Build a single text for detect_headings_and_map
        combined = "\n\n".join([p["text"] for p in merged_paras])
        para_heading_map = detect_headings_and_map(combined)
        # map by index
        idx_to_heading = {m["para_index"]: m["heading"] for m in para_heading_map}
        for idx in range(len(merged_paras)):
            headings_map.append(idx_to_heading.get(idx))
    else:
        headings_map = [None] * len(merged_paras)

    # Build chunks
    chunks: List[Dict] = []
    chunk_buf: List[str] = []
    chunk_pages: List[int] = []
    chunk_headings: List[Optional[str]] = []
    current_len = 0
    para_index = 0
    for idx, para in enumerate(merged_paras):
        para_text = para["text"].strip()
        para_page = para.get("page")
        if para_text == "":
            continue
        # If adding this paragraph would exceed max_chunk_chars and chunk_buf not empty, flush chunk
        if current_len + len(para_text) > max_chunk_chars and chunk_buf:
            chunk_text = "\n\n".join(chunk_buf).strip()
            page_start = next((p for p in chunk_pages if p is not None), None)
            page_end = next((p for p in reversed(chunk_pages) if p is not None), None)
            heading = chunk_headings[0] if chunk_headings else None
            chunk_id = len(chunks)
            chunks.append({
                "chunk_id": chunk_id,
                "text": chunk_text,
                "heading": heading,
                "page_start": page_start,
                "page_end": page_end,
            })
            # prepare overlap: keep last overlap_chars characters
            overlap_text = chunk_text[-overlap_chars:] if overlap_chars and len(chunk_text) > overlap_chars else chunk_text
            chunk_buf = [overlap_text] if overlap_text else []
            chunk_pages = []
            chunk_headings = []
            current_len = len(overlap_text)

        chunk_buf.append(para_text)
        if para_page is not None:
            chunk_pages.append(para_page)
        chunk_headings.append(headings_map[idx] if idx < len(headings_map) else None)
        current_len += len(para_text)

    # flush final chunk
    if chunk_buf:
        chunk_text = "\n\n".join(chunk_buf).strip()
        page_start = next((p for p in chunk_pages if p is not None), None)
        page_end = next((p for p in reversed(chunk_pages) if p is not None), None)
        heading = chunk_headings[0] if chunk_headings else None
        chunk_id = len(chunks)
        chunks.append({
            "chunk_id": chunk_id,
            "text": chunk_text,
            "heading": heading,
            "page_start": page_start,
            "page_end": page_end,
        })

    # Final trimming: remove leading/trailing punctuation from texts
    for c in chunks:
        c["text"] = c["text"].strip(" \n\t-:;,.()[]{}").strip()

    return chunks


if __name__ == "__main__":
    import argparse
    import os

    logging.basicConfig(level=logging.INFO)
    parser = argparse.ArgumentParser(description="Test file_service extraction and chunking")
    parser.add_argument("--file", type=str, help="Path to sample PDF/DOCX/PPTX/TXT")
    parser.add_argument("--max_chunk", type=int, default=800)
    args = parser.parse_args()
    sample = args.file or os.environ.get("SAMPLE_DOC")
    if not sample:
        print("Provide a sample file via --file or SAMPLE_DOC env var to run the quick test.")
        raise SystemExit(1)

    try:
        chunks = extract_text_chunks(sample, max_chunk_chars=args.max_chunk)
        print(f"Extracted {len(chunks)} chunks from {sample}")
        if chunks:
            print("--- First chunk preview ---")
            print(chunks[0]["text"][:200])
            print("Heading:", chunks[0]["heading"]) 
    except Exception as e:
        logger.exception("Test extraction failed: %s", e)
        raise
