import io
import re
from typing import Optional
import pypdf
import docx
from PIL import Image
from pptx import Presentation
from striprtf.striprtf import rtf_to_text
import markdown
from bs4 import BeautifulSoup

class ParserService:
    @staticmethod
    def extract_text(file_contents: bytes, filename: str) -> str:
        """Extract text from various file formats with preprocessing"""
        filename_lower = filename.lower()
        
        try:
            if filename_lower.endswith(".pdf"):
                text = ParserService._parse_pdf(file_contents)
            elif filename_lower.endswith((".docx", ".doc")):
                text = ParserService._parse_docx(file_contents)
            elif filename_lower.endswith(".txt"):
                text = file_contents.decode("utf-8", errors='ignore')
            elif filename_lower.endswith((".pptx", ".ppt")):
                text = ParserService._parse_pptx(file_contents)
            elif filename_lower.endswith(".rtf"):
                text = ParserService._parse_rtf(file_contents)
            elif filename_lower.endswith((".md", ".markdown")):
                text = ParserService._parse_markdown(file_contents)
            elif filename_lower.endswith((".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff")):
                text = ParserService._parse_image(file_contents)
            else:
                # Try to decode as text
                text = file_contents.decode("utf-8", errors='ignore')
            
            # Preprocess and clean the text
            return ParserService._preprocess_text(text)
            
        except Exception as e:
            print(f"Error parsing {filename}: {e}")
            # Fallback: try to decode as text
            try:
                return ParserService._preprocess_text(file_contents.decode("utf-8", errors='ignore'))
            except:
                return "Error: Could not extract text from this file."

    @staticmethod
    def _parse_pdf(file_contents: bytes) -> str:
        text = ""
        try:
            reader = pypdf.PdfReader(io.BytesIO(file_contents))
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
        except Exception as e:
            print(f"PDF parsing error: {e}")
        return text

    @staticmethod
    def _parse_docx(file_contents: bytes) -> str:
        text = ""
        try:
            doc = docx.Document(io.BytesIO(file_contents))
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text += cell.text + " "
                    text += "\n"
                    
        except Exception as e:
            print(f"DOCX parsing error: {e}")
        return text

    @staticmethod
    def _parse_pptx(file_contents: bytes) -> str:
        text = ""
        try:
            prs = Presentation(io.BytesIO(file_contents))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
        except Exception as e:
            print(f"PPTX parsing error: {e}")
        return text

    @staticmethod
    def _parse_rtf(file_contents: bytes) -> str:
        try:
            rtf_content = file_contents.decode("utf-8", errors='ignore')
            return rtf_to_text(rtf_content)
        except Exception as e:
            print(f"RTF parsing error: {e}")
            return ""

    @staticmethod
    def _parse_markdown(file_contents: bytes) -> str:
        try:
            md_content = file_contents.decode("utf-8", errors='ignore')
            html = markdown.markdown(md_content)
            # Convert HTML to plain text
            soup = BeautifulSoup(html, 'html.parser')
            return soup.get_text()
        except Exception as e:
            print(f"Markdown parsing error: {e}")
            return file_contents.decode("utf-8", errors='ignore')

    @staticmethod
    def _parse_image(file_contents: bytes) -> str:
        """Extract text from images using OCR (requires pytesseract)"""
        try:
            # For now, return a placeholder since OCR setup can be complex
            # In production, you would use pytesseract here
            return "Image content detected. OCR functionality requires additional setup."
        except Exception as e:
            print(f"Image parsing error: {e}")
            return "Could not extract text from image."

    @staticmethod
    def _preprocess_text(text: str) -> str:
        """Clean and preprocess extracted text for better AI processing"""
        if not text or not text.strip():
            return ""
        
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove special characters that might confuse the model
        text = re.sub(r'[^\w\s\.\,\?\!\:\;\-\(\)\[\]\{\}\"\'\/]', ' ', text)
        
        # Fix common OCR/extraction errors
        text = re.sub(r'\b([a-z])([A-Z])', r'\1 \2', text)  # Split camelCase
        text = re.sub(r'([a-z])(\d)', r'\1 \2', text)  # Separate letters from numbers
        text = re.sub(r'(\d)([a-z])', r'\1 \2', text)  # Separate numbers from letters
        
        # Remove lines that are too short (likely noise)
        lines = text.split('\n')
        meaningful_lines = []
        for line in lines:
            line = line.strip()
            if len(line) > 10:  # Only keep lines with substantial content
                meaningful_lines.append(line)
        
        # Join lines back together
        text = ' '.join(meaningful_lines)
        
        # Remove excessive punctuation
        text = re.sub(r'[\.]{3,}', '...', text)
        text = re.sub(r'[-]{3,}', '---', text)
        
        # Ensure sentences end properly
        text = re.sub(r'([a-zA-Z])([A-Z])', r'\1. \2', text)
        
        # Final cleanup
        text = re.sub(r'\s+', ' ', text).strip()
        
        return text

    @staticmethod
    def chunk_text(text: str, chunk_size: int = 800, overlap: int = 100) -> list[str]:
        """Split text into overlapping chunks for better context preservation"""
        if len(text) <= chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # Try to break at sentence boundaries
            if end < len(text):
                # Look for sentence endings within the last 100 characters
                sentence_end = text.rfind('.', start, end)
                if sentence_end > start + chunk_size - 200:
                    end = sentence_end + 1
                else:
                    # Look for word boundaries
                    word_end = text.rfind(' ', start, end)
                    if word_end > start + chunk_size - 100:
                        end = word_end
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            # Move start position with overlap
            start = end - overlap
            if start >= len(text):
                break
        
        return chunks